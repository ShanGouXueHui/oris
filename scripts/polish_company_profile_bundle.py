#!/usr/bin/env python3
import argparse
import json
import re
from copy import deepcopy
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parents[1]
CFG_PATH = ROOT / "config" / "polish_render_config.json"


def load_json(path: Path):
    return json.loads(path.read_text(encoding="utf-8"))


def dump_json(path: Path, data):
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(data, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def cfg():
    try:
        return load_json(CFG_PATH)
    except Exception:
        return {
            "version": 1,
            "default_theme": "executive_blue",
            "themes": {
                "executive_blue": {
                    "bg": "#F7F9FC",
                    "title": "#0F2747",
                    "text": "#1F2937",
                    "muted": "#6B7280",
                    "accent": "#2563EB",
                    "accent2": "#0EA5E9",
                    "card": "#FFFFFF",
                    "line": "#D1D5DB",
                    "risk": "#DC2626",
                    "ok": "#059669"
                }
            },
            "profile_theme_map": {},
            "unit_display": {},
            "profile_metric_priority": {
                "generic_company": ["revenue", "gross_profit", "operating_profit", "free_cash_flow"]
            },
            "metric_scoring_bonus": {},
            "metric_candidate_rules": {},
            "internal_line_patterns": [],
            "ppt_layout": {
                "slide_width_inches": 13.333,
                "slide_height_inches": 7.5,
                "max_exec_summary_items": 4,
                "max_dashboard_items": 8,
                "max_two_col_items_per_page": 12,
                "max_risk_items": 4,
                "max_tracking_items": 8
            }
        }


def normalize_text(text: str) -> str:
    s = str(text or "")
    s = s.replace("\x00", " ").replace("\ufeff", " ").replace("\r", "\n")
    s = re.sub(r"[ \t]+", " ", s)
    s = re.sub(r"\n{3,}", "\n\n", s)
    return s.strip()


def first_existing_path(path_str: str):
    if not path_str:
        return None
    p = Path(path_str)
    if p.is_absolute() and p.exists():
        return p
    rp = ROOT / path_str
    if rp.exists():
        return rp
    return None


def parse_date(s: str):
    s = str(s or "").strip()
    if not s:
        return datetime(1970, 1, 1)
    try:
        return datetime.fromisoformat(s.replace("Z", "+00:00"))
    except Exception:
        return datetime(1970, 1, 1)


def pretty_number(v: float, decimals: int = 1) -> str:
    if abs(v - round(v)) < 1e-9:
        return str(int(round(v)))
    return f"{v:.{decimals}f}"


def rgb(hex_str: str) -> RGBColor:
    s = (hex_str or "#000000").strip().lstrip("#")
    if len(s) != 6:
        s = "000000"
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


def theme_for_profile(focus_profile: str):
    c = cfg()
    name = (c.get("profile_theme_map") or {}).get(focus_profile) or c.get("default_theme") or "executive_blue"
    theme = (c.get("themes") or {}).get(name) or (c.get("themes") or {}).get(c.get("default_theme") or "") or {}
    return {
        "name": name,
        "palette": theme
    }


def add_text(slide, left, top, width, height, text, size=18, bold=False, color=None, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if hasattr(p, "clear"):
        p.clear()
    p.alignment = align
    run = p.add_run()
    run.text = str(text) if str(text or "").strip() else " "
    run.font.size = Pt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    return box


def add_rect(slide, left, top, width, height, fill_hex, line_hex=None, rounded=True):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill_hex)
    shp.line.color.rgb = rgb(line_hex or fill_hex)
    return shp


def apply_bg(slide, fill_hex):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(fill_hex)


def unit_to_cn(value, unit: str) -> str:
    c = cfg()
    rules = c.get("unit_display") or {}
    rule = rules.get(str(unit or "").strip())
    try:
        v = float(value)
    except Exception:
        return f"{value}{unit or ''}".strip()

    if not rule:
        return f"{pretty_number(v)}{unit or ''}".strip()

    decimals = int(rule.get("decimals", 1))
    kind = rule.get("kind") or "suffix"
    prefix = str(rule.get("prefix", ""))
    suffix = str(rule.get("suffix", ""))
    multiplier = float(rule.get("multiplier", 1))
    divisor = float(rule.get("divisor", 1))

    vv = v
    if kind in {"multiply_suffix", "multiply_prefix_suffix"}:
        vv = v * multiplier
    elif kind in {"divide_suffix", "divide_prefix_suffix"}:
        vv = v / divisor

    if kind in {"suffix", "multiply_suffix", "divide_suffix"}:
        return f"{pretty_number(vv, decimals)}{suffix}".strip()
    if kind in {"prefix_suffix", "multiply_prefix_suffix", "divide_prefix_suffix"}:
        return f"{prefix}{pretty_number(vv, decimals)}{suffix}".strip()

    return f"{pretty_number(vv, decimals)}{suffix}".strip()


def clean_line(line: str) -> str:
    t = normalize_text(line)
    t = re.sub(r"\s+Yo Y\b", " YoY", t, flags=re.I)
    t = re.sub(r"\bY o Y\b", "YoY", t, flags=re.I)
    t = re.sub(r"\bQ o Q\b", "QoQ", t, flags=re.I)
    t = t.strip("•▪●◦- ")
    return t


def is_internal_line(line: str) -> bool:
    pats = [str(x).lower() for x in (cfg().get("internal_line_patterns") or [])]
    t = normalize_text(line).lower()
    return any(x in t for x in pats)


def clean_section_items(items):
    out = []
    seen = set()
    for x in items or []:
        t = clean_line(x)
        if not t:
            continue
        if is_internal_line(t):
            continue
        key = t.lower()
        if key in seen:
            continue
        seen.add(key)
        out.append(t)
    return out


def priority_codes(focus_profile: str):
    c = cfg()
    p = c.get("profile_metric_priority") or {}
    return p.get(focus_profile) or p.get("generic_company") or []


def metric_candidate_score(row: dict, focus_profile: str) -> int:
    c = cfg()
    code = str(row.get("metric_code") or "")
    unit = str(row.get("metric_unit") or "")
    try:
        value = float(row.get("metric_value"))
    except Exception:
        value = 0.0

    score = 0
    score += int((c.get("metric_scoring_bonus") or {}).get(unit, 0))

    special = (c.get("metric_candidate_rules") or {}).get(f"{code}|{unit}") or {}
    if special:
        ok = True
        if "min" in special and value < float(special["min"]):
            ok = False
        if "max" in special and value > float(special["max"]):
            ok = False
        if ok:
            score += int(special.get("bonus", 0))

    score += int(parse_date(row.get("observation_date")).timestamp() // 86400)
    return score


def line_from_metric(row: dict) -> str:
    name = str(row.get("metric_name") or row.get("metric_code") or "").strip()
    if not name:
        return ""
    value_text = unit_to_cn(row.get("metric_value"), row.get("metric_unit"))
    obs = str(row.get("observation_date") or "").strip()
    return f"{name}: {value_text}（{obs}）" if obs else f"{name}: {value_text}"


def top_structured_metrics(profile: dict, focus_profile: str, limit: int = 8):
    rows = profile.get("recent_metric_observations") or []
    preferred = priority_codes(focus_profile)
    if not rows or not preferred:
        return []

    rank = {code: idx for idx, code in enumerate(preferred)}
    best_by_code = {}

    for row in rows:
        code = str(row.get("metric_code") or "").strip()
        if not code or code not in rank:
            continue
        cand = deepcopy(row)
        cand["_score"] = metric_candidate_score(cand, focus_profile)
        old = best_by_code.get(code)
        if old is None or cand["_score"] > old["_score"]:
            best_by_code[code] = cand

    out = []
    for code in preferred:
        row = best_by_code.get(code)
        if not row:
            continue
        line = line_from_metric(row)
        if line:
            out.append(line)
        if len(out) >= limit:
            break
    return out


def summarize_top_metrics(lines, limit=3):
    vals = [clean_line(x) for x in (lines or []) if clean_line(x)]
    return vals[:limit]


def enrich_sections(bundle: dict, profile: dict, focus_profile: str):
    c = cfg()
    sections = deepcopy(bundle.get("sections") or {})

    for key in ["company_view", "technology_view", "industry_view", "risk_view", "tracking_kpis"]:
        sections[key] = clean_section_items(sections.get(key) or [])

    structured = top_structured_metrics(
        profile,
        focus_profile,
        limit=int((c.get("ppt_layout") or {}).get("max_tracking_items", 8))
    )
    if structured:
        sections["tracking_kpis"] = structured

    exec_summary = clean_section_items(sections.get("executive_summary") or [])
    if len(exec_summary) < 3:
        exec_summary = [
            "已完成基于当前可得证据的公司洞察与结构化整理。",
            f"当前焦点画像：{focus_profile}。",
            "建议以业务结构、核心产品/技术、竞争位置、风险与跟踪指标五层持续跟踪。"
        ]

    tops = summarize_top_metrics(sections.get("tracking_kpis") or [], limit=3)
    if tops:
        line4 = "优先关注指标：" + "；".join(tops)
        if len(exec_summary) >= 4:
            exec_summary[3] = line4
        else:
            exec_summary.append(line4)

    sections["executive_summary"] = exec_summary[:int((c.get("ppt_layout") or {}).get("max_exec_summary_items", 4))]
    return sections


def make_polished_ppt(bundle: dict, out_path: Path):
    c = cfg()
    sections = bundle.get("sections") or {}
    focus_profile = bundle.get("focus_profile") or "generic_company"
    company_name = ((bundle.get("request") or {}).get("company_name") or "Company").strip()
    theme = theme_for_profile(focus_profile)["palette"]
    layout = c.get("ppt_layout") or {}

    prs = Presentation()
    prs.slide_width = Inches(float(layout.get("slide_width_inches", 13.333)))
    prs.slide_height = Inches(float(layout.get("slide_height_inches", 7.5)))

    # cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_bg(slide, theme["bg"])
    add_rect(slide, Inches(0.55), Inches(0.65), Inches(0.18), Inches(5.9), theme["accent"], theme["accent"])
    add_text(slide, Inches(0.95), Inches(1.0), Inches(10.8), Inches(0.9),
             f"{company_name} 公司洞察", size=28, bold=True, color=rgb(theme["title"]))
    add_text(slide, Inches(0.98), Inches(1.85), Inches(8.0), Inches(0.4),
             f"焦点画像：{focus_profile}", size=13, color=rgb(theme["muted"]))
    add_text(slide, Inches(0.98), Inches(2.35), Inches(10.8), Inches(1.4),
             "结构：执行摘要｜关键指标｜公司与商业模式｜技术与产品｜行业与竞争｜风险与跟踪",
             size=16, color=rgb(theme["text"]))
    add_rect(slide, Inches(0.98), Inches(3.4), Inches(3.0), Inches(0.08), theme["accent2"], theme["accent2"])
    add_text(slide, Inches(0.98), Inches(5.9), Inches(8.0), Inches(0.3),
             "ORIS Polished Deck v1", size=11, color=rgb(theme["muted"]))

    # exec summary
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_bg(slide, theme["bg"])
    add_text(slide, Inches(0.6), Inches(0.4), Inches(6), Inches(0.5),
             "执行摘要", size=24, bold=True, color=rgb(theme["title"]))
    items = (sections.get("executive_summary") or [])[:int(layout.get("max_exec_summary_items", 4))]
    positions = [(0.7, 1.2), (3.55, 1.2), (6.4, 1.2), (9.25, 1.2)]
    card_titles = ["核心判断", "画像定位", "推进逻辑", "优先指标"]
    for i, (x, y) in enumerate(positions[:len(items)]):
        add_rect(slide, Inches(x), Inches(y), Inches(2.55), Inches(2.7), theme["card"], theme["line"])
        add_rect(slide, Inches(x), Inches(y), Inches(2.55), Inches(0.12), theme["accent"], theme["accent"])
        add_text(slide, Inches(x + 0.15), Inches(y + 0.25), Inches(2.1), Inches(0.3),
                 card_titles[i], size=12, bold=True, color=rgb(theme["accent"]))
        add_text(slide, Inches(x + 0.15), Inches(y + 0.62), Inches(2.2), Inches(1.85),
                 items[i], size=13, color=rgb(theme["text"]))

    # KPI dashboard
    kpis = (sections.get("tracking_kpis") or [])[:int(layout.get("max_dashboard_items", 8))]
    if kpis:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        apply_bg(slide, theme["bg"])
        add_text(slide, Inches(0.6), Inches(0.4), Inches(6), Inches(0.5),
                 "关键指标仪表板", size=24, bold=True, color=rgb(theme["title"]))
        for idx, item in enumerate(kpis):
            row = idx // 4
            col = idx % 4
            x = 0.7 + col * 3.05
            y = 1.2 + row * 2.35
            add_rect(slide, Inches(x), Inches(y), Inches(2.7), Inches(1.9), theme["card"], theme["line"])
            add_rect(slide, Inches(x), Inches(y), Inches(2.7), Inches(0.1), theme["accent2"], theme["accent2"])
            add_text(slide, Inches(x + 0.15), Inches(y + 0.22), Inches(2.35), Inches(1.35),
                     item, size=12, color=rgb(theme["text"]))

    def two_col_page(title, items):
        if not items:
            return
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        apply_bg(slide, theme["bg"])
        add_text(slide, Inches(0.6), Inches(0.4), Inches(8), Inches(0.5),
                 title, size=24, bold=True, color=rgb(theme["title"]))
        max_items = int(layout.get("max_two_col_items_per_page", 12))
        items = items[:max_items]
        cols = [items[::2], items[1::2]]
        for ci, col in enumerate(cols):
            x = 0.7 if ci == 0 else 6.55
            add_rect(slide, Inches(x), Inches(1.15), Inches(5.95), Inches(5.7), theme["card"], theme["line"])
            y = 1.38
            for item in col[:6]:
                add_text(slide, Inches(x + 0.18), Inches(y), Inches(5.45), Inches(0.72),
                         "• " + item, size=13, color=rgb(theme["text"]))
                y += 0.84

    two_col_page("公司与商业模式", sections.get("company_view") or [])
    two_col_page("技术与产品能力", sections.get("technology_view") or [])
    two_col_page("行业与竞争", sections.get("industry_view") or [])

    risks = (sections.get("risk_view") or [])[:int(layout.get("max_risk_items", 4))]
    if risks:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        apply_bg(slide, theme["bg"])
        add_text(slide, Inches(0.6), Inches(0.4), Inches(6), Inches(0.5),
                 "风险矩阵", size=24, bold=True, color=rgb(theme["title"]))
        for idx, item in enumerate(risks):
            x = 0.8 + (idx % 2) * 6.0
            y = 1.3 + (idx // 2) * 2.35
            add_rect(slide, Inches(x), Inches(y), Inches(5.4), Inches(1.7), theme["card"], theme["risk"])
            add_text(slide, Inches(x + 0.18), Inches(y + 0.18), Inches(5.0), Inches(1.1),
                     item, size=13, color=rgb(theme["text"]))

    if kpis:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        apply_bg(slide, theme["bg"])
        add_text(slide, Inches(0.6), Inches(0.4), Inches(6), Inches(0.5),
                 "持续跟踪清单", size=24, bold=True, color=rgb(theme["title"]))
        y = 1.15
        for item in kpis[:int(layout.get("max_tracking_items", 8))]:
            add_rect(slide, Inches(0.8), Inches(y), Inches(11.2), Inches(0.58), theme["card"], theme["line"])
            add_rect(slide, Inches(0.8), Inches(y), Inches(0.12), Inches(0.58), theme["accent"], theme["accent"])
            add_text(slide, Inches(1.0), Inches(y + 0.1), Inches(10.7), Inches(0.3),
                     item, size=12, color=rgb(theme["text"]))
            y += 0.7

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--bundle-json-path", required=True)
    ap.add_argument("--output-json-path", default="")
    ap.add_argument("--output-pptx-path", default="")
    args = ap.parse_args()

    bundle_path = first_existing_path(args.bundle_json_path)
    if not bundle_path:
        raise SystemExit("bundle json not found")

    bundle = load_json(bundle_path)
    focus_profile = bundle.get("focus_profile") or ((bundle.get("request") or {}).get("focus_profile")) or "generic_company"

    input_json_path = (bundle.get("request") or {}).get("input_json_path") or ""
    profile_json_path = first_existing_path(input_json_path)
    profile_data = load_json(profile_json_path) if profile_json_path and profile_json_path.exists() else {}
    profile = profile_data.get("company_profile") or profile_data.get("db_backed_profile") or {}

    polished = deepcopy(bundle)
    polished["sections"] = enrich_sections(bundle, profile, focus_profile)
    polished["postprocess_mode"] = "polished_v2_config_driven"
    polished["postprocess_ts"] = datetime.utcnow().isoformat() + "Z"
    polished["postprocess_config_path"] = str(CFG_PATH.relative_to(ROOT))

    out_json = Path(args.output_json_path) if args.output_json_path else bundle_path.with_name("company_profile_bundle.polished.json")
    if not out_json.is_absolute():
        out_json = (bundle_path.parent / out_json.name) if out_json.name == out_json.as_posix() else (ROOT / out_json)

    out_pptx = Path(args.output_pptx_path) if args.output_pptx_path else bundle_path.with_name("company_profile_deck.polished.pptx")
    if not out_pptx.is_absolute():
        out_pptx = (bundle_path.parent / out_pptx.name) if out_pptx.name == out_pptx.as_posix() else (ROOT / out_pptx)

    dump_json(out_json, polished)
    make_polished_ppt(polished, out_pptx)

    print(json.dumps({
        "ok": True,
        "bundle_json_path": str(bundle_path.relative_to(ROOT)) if bundle_path.is_relative_to(ROOT) else str(bundle_path),
        "output_json_path": str(out_json.relative_to(ROOT)) if out_json.is_relative_to(ROOT) else str(out_json),
        "output_pptx_path": str(out_pptx.relative_to(ROOT)) if out_pptx.is_relative_to(ROOT) else str(out_pptx),
        "focus_profile": focus_profile,
        "postprocess_config_path": str(CFG_PATH.relative_to(ROOT)),
        "tracking_kpi_count": len((polished.get("sections") or {}).get("tracking_kpis") or []),
        "exec_summary_preview": ((polished.get("sections") or {}).get("executive_summary") or [])[:4]
    }, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
