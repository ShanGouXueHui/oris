#!/usr/bin/env python3
import argparse
import json
from collections import defaultdict
from datetime import datetime, timezone
from pathlib import Path

from docx import Document
from docx.shared import Pt
from openpyxl import Workbook
from pptx import Presentation

import sys
ROOT = Path(__file__).resolve().parents[2]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from scripts.lib.oris_llm_client import call_oris_text

ROOT = Path(__file__).resolve().parents[2]

def utc_now():
    return datetime.now(timezone.utc).isoformat()

def ts_compact():
    return datetime.now(timezone.utc).strftime("%Y%m%d_%H%M%S")

def load_json(path: Path):
    return json.loads(path.read_text(encoding="utf-8"))

def load_input_payload(value: str):
    s = (value or "").strip()
    if not s:
        raise ValueError("empty input json path/payload")
    if s.startswith("{"):
        return json.loads(s)
    p = Path(s)
    return json.loads(p.read_text(encoding="utf-8"))

def resolve_case_payload(value: str):
    payload = load_input_payload(value)

    # 1) already final case json
    if isinstance(payload, dict) and (
        payload.get("case_code") or ((payload.get("request") or {}).get("case_code"))
    ):
        return payload

    # 2) wrapper json from pipeline/report request
    if isinstance(payload, dict) and payload.get("input_json_path"):
        inner = payload.get("input_json_path")
        if not inner:
            raise ValueError("input_json_path is empty in wrapper payload")
        case = load_input_payload(inner)
        if not isinstance(case, dict) or not (
            case.get("case_code") or ((case.get("request") or {}).get("case_code"))
        ):
            raise ValueError("wrapped input_json_path did not resolve to account strategy case json")
        return case

    raise ValueError("unsupported account strategy input payload shape")

def write_json(path: Path, data: dict):
    path.write_text(json.dumps(data, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")

def normalize(v):
    if v is None:
        return ""
    if isinstance(v, (dict, list)):
        return json.dumps(v, ensure_ascii=False)
    return str(v)

def collect_nested_rows(obj, target_key, out):
    if isinstance(obj, dict):
        for k, v in obj.items():
            if k == target_key and isinstance(v, list):
                out.extend(v)
            else:
                collect_nested_rows(v, target_key, out)
    elif isinstance(obj, list):
        for item in obj:
            collect_nested_rows(item, target_key, out)

def safe_json_from_text(text: str):
    text = (text or "").strip()
    if not text:
        return None
    try:
        return json.loads(text)
    except Exception:
        pass
    import re
    m = re.search(r'(\{.*\})', text, flags=re.S)
    if m:
        try:
            return json.loads(m.group(1))
        except Exception:
            return None
    return None

def build_synthesis(case: dict, evidence_rows: list, citation_rows: list):
    prompt = {
        "task": "生成高质量商业洞察内容",
        "requirements": {
            "style": "专业、分层、可用于商务汇报，不允许口水化",
            "must_cover": [
                "行业层",
                "技术栈层",
                "客户场景层",
                "竞争格局层",
                "合作方案层",
                "风险与下一步"
            ],
            "must_distinguish": ["事实", "推断", "建议", "风险"],
            "must_use_evidence": True
        },
        "input": {
            "case_code": case.get("case_code"),
            "questions": case.get("questions") or [],
            "frameworks": case.get("frameworks") or [],
            "report_sections": case.get("report_sections") or [],
            "entity_summaries": case.get("entity_summaries") or [],
            "competitor_matrix": case.get("competitor_matrix") or [],
            "capability_mapping": case.get("capability_mapping") or [],
            "analysis_sections": case.get("analysis_sections") or [],
            "recommendation_titles": case.get("recommendation_titles") or [],
            "evidence_sample": evidence_rows[:40],
            "citation_sample": citation_rows[:40]
        },
        "output_schema": {
            "executive_summary": ["..."],
            "industry_view": ["..."],
            "technology_stack_view": ["..."],
            "customer_scene_view": {
                "引望": ["..."],
                "北汽": ["..."]
            },
            "competitive_view": ["..."],
            "joint_solution_recommendations": ["..."],
            "risks_and_next_steps": ["..."]
        }
    }

    try:
        resp = call_oris_text(
            "请基于下面JSON生成严格JSON，不要输出解释文字：\n" + json.dumps(prompt, ensure_ascii=False),
            role="free_fallback",
            timeout_seconds=180
        )
        if resp.get("ok"):
            parsed = safe_json_from_text(resp.get("text", ""))
            if parsed:
                return {"mode": "llm", "sections": parsed}
    except Exception:
        pass

    # fallback deterministic
    entity_summaries = case.get("entity_summaries") or []
    competitor_matrix = case.get("competitor_matrix") or []
    capability_mapping = case.get("capability_mapping") or []

    sections = {
        "executive_summary": [
            "本报告围绕伙伴、云平台、客户、竞争对手四层关系图展开，目标是形成可落地的联合方案而非泛泛介绍。",
            f"当前证据条数={len(evidence_rows)}，引用条数={len(citation_rows)}，可支撑正式商务交流材料。"
        ],
        "industry_view": [
            "汽车行业已从单点智能功能竞争，转向软件定义汽车、车云闭环、数据驱动迭代与全球化交付能力竞争。",
            "对工程服务伙伴而言，价值不只在研发外包，而在把 AI、验证、数据、云平台、合规与全球交付打通。"
        ],
        "technology_stack_view": [
            "建议按五层技术栈展开：算力与云底座、数据与MLOps、模型与AI开发平台、车端/云端工程与验证、运营闭环与质量分析。",
            "华为云的价值不应只写成‘有AI能力’，而要展开到 ModelArts/MLOps/数据处理/推理部署/车云协同等能力层。"
        ],
        "customer_scene_view": {
            "引望": [
                "引望更偏智能汽车平台/方案供给方，重点应放在平台能力、生态适配、开发者体系、量产交付效率。",
                "联合方案应强调 AI 工程效率、工具链、验证闭环与生态输出能力。"
            ],
            "北汽": [
                "北汽更偏整车OEM，重点应放在车型项目落地、供应链协同、成本、上市节奏、海外车型复制与品牌差异化。",
                "联合方案应强调车企数字化研发、质量闭环、营销/售后数据智能与海外拓展支撑。"
            ]
        },
        "competitive_view": [
            "竞争对手比较不宜只讲‘谁强谁弱’，应拆成工程深度、AI能力、全栈方案、生态、交付成本、区域覆盖六个维度。",
            f"当前竞争矩阵条目={len(competitor_matrix)}，建议在正式汇报中把原始证据与原始出处直接挂到每个维度。"
        ],
        "joint_solution_recommendations": [
            x.get("joint_value") for x in capability_mapping[:6] if x.get("joint_value")
        ] or ["建议把联合方案拆成 SDV/智能驾驶工程工厂、车云运营闭环、海外本地化交付 三条主线。"],
        "risks_and_next_steps": [
            "对外材料必须避免抽象打分，优先展示原始证据、出处、能力映射和客户场景匹配。",
            "下一步应继续补年报/投资者关系/权威新闻/第三方研究源，增强规模、财务、客户案例与生态合作的数据密度。"
        ]
    }
    return {"mode": "deterministic_fallback", "sections": sections}

def add_bullets(slide, title, items):
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    items = items or ["N/A"]
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = str(item)

def make_docx(path: Path, case: dict, sections: dict, evidence_rows: list, citation_rows: list):
    doc = Document()
    h = doc.add_heading("Akkodis × Huawei Cloud 商业洞察报告", 0)
    h.runs[0].font.size = Pt(20)

    doc.add_heading("一、执行摘要", level=1)
    for x in sections.get("executive_summary") or []:
        doc.add_paragraph(str(x), style="List Bullet")

    doc.add_heading("二、行业层", level=1)
    for x in sections.get("industry_view") or []:
        doc.add_paragraph(str(x), style="List Bullet")

    doc.add_heading("三、AI技术栈层", level=1)
    for x in sections.get("technology_stack_view") or []:
        doc.add_paragraph(str(x), style="List Bullet")

    doc.add_heading("四、客户场景层", level=1)
    csv = sections.get("customer_scene_view") or {}
    for k, items in csv.items():
        doc.add_heading(str(k), level=2)
        for x in items or []:
            doc.add_paragraph(str(x), style="List Bullet")

    doc.add_heading("五、竞争格局层", level=1)
    for x in sections.get("competitive_view") or []:
        doc.add_paragraph(str(x), style="List Bullet")

    doc.add_heading("六、联合方案建议", level=1)
    for x in sections.get("joint_solution_recommendations") or []:
        doc.add_paragraph(str(x), style="List Bullet")

    doc.add_heading("七、风险与下一步", level=1)
    for x in sections.get("risks_and_next_steps") or []:
        doc.add_paragraph(str(x), style="List Bullet")

    doc.add_heading("八、实体摘要", level=1)
    table = doc.add_table(rows=1, cols=6)
    hdr = table.rows[0].cells
    hdr[0].text = "Entity"
    hdr[1].text = "Type"
    hdr[2].text = "Snapshots"
    hdr[3].text = "Evidence"
    hdr[4].text = "Citations"
    hdr[5].text = "Notes"
    for row in case.get("entity_summaries") or []:
        cells = table.add_row().cells
        cells[0].text = normalize(row.get("entity_name"))
        cells[1].text = normalize(row.get("entity_type"))
        cells[2].text = normalize(row.get("source_snapshot_count_recent"))
        cells[3].text = normalize(row.get("evidence_item_count_recent"))
        cells[4].text = normalize(row.get("citation_count_recent"))
        cells[5].text = normalize((row.get("sample_metric_codes") or [])[:3])

    doc.add_heading("九、竞争矩阵（原始量）", level=1)
    table = doc.add_table(rows=1, cols=6)
    hdr = table.rows[0].cells
    hdr[0].text = "Entity"
    hdr[1].text = "Type"
    hdr[2].text = "Snapshots"
    hdr[3].text = "Evidence"
    hdr[4].text = "Metrics"
    hdr[5].text = "Citations"
    for row in case.get("competitor_matrix") or []:
        cells = table.add_row().cells
        cells[0].text = normalize(row.get("entity_name"))
        cells[1].text = normalize(row.get("entity_type"))
        cells[2].text = normalize(row.get("source_snapshot_count_recent"))
        cells[3].text = normalize(row.get("evidence_item_count_recent"))
        cells[4].text = normalize(row.get("metric_observation_count_recent"))
        cells[5].text = normalize(row.get("citation_count_recent"))

    doc.add_heading("十、证据样本", level=1)
    for row in evidence_rows[:30]:
        doc.add_paragraph(
            f"{normalize(row.get('evidence_title'))}\n{normalize(row.get('evidence_text'))}",
            style="List Bullet"
        )

    doc.add_heading("十一、引用附录", level=1)
    for row in citation_rows[:80]:
        doc.add_paragraph(
            f"{normalize(row.get('citation_label'))}\n{normalize(row.get('citation_url'))}",
            style="List Bullet"
        )

    doc.save(path)

def make_xlsx(path: Path, case: dict, evidence_rows: list, citation_rows: list, snapshot_rows: list, metric_rows: list):
    wb = Workbook()

    ws = wb.active
    ws.title = "entity_summaries"
    ws.append(["entity_name", "entity_type", "source_snapshot_count_recent", "evidence_item_count_recent", "metric_observation_count_recent", "citation_count_recent"])
    for row in case.get("entity_summaries") or []:
        ws.append([
            normalize(row.get("entity_name")),
            normalize(row.get("entity_type")),
            normalize(row.get("source_snapshot_count_recent")),
            normalize(row.get("evidence_item_count_recent")),
            normalize(row.get("metric_observation_count_recent")),
            normalize(row.get("citation_count_recent"))
        ])

    ws2 = wb.create_sheet("competitor_matrix")
    ws2.append(["entity_name", "entity_type", "dimensions", "source_snapshot_count_recent", "evidence_item_count_recent", "metric_observation_count_recent", "citation_count_recent"])
    for row in case.get("competitor_matrix") or []:
        ws2.append([
            normalize(row.get("entity_name")),
            normalize(row.get("entity_type")),
            normalize((row.get("dimensions") or [])),
            normalize(row.get("source_snapshot_count_recent")),
            normalize(row.get("evidence_item_count_recent")),
            normalize(row.get("metric_observation_count_recent")),
            normalize(row.get("citation_count_recent"))
        ])

    ws3 = wb.create_sheet("evidence_raw")
    ws3.append(["id", "company_id", "source_snapshot_id", "evidence_type", "evidence_title", "evidence_text", "evidence_date", "confidence_score"])
    for row in evidence_rows:
        ws3.append([
            normalize(row.get("id")),
            normalize(row.get("company_id")),
            normalize(row.get("source_snapshot_id")),
            normalize(row.get("evidence_type")),
            normalize(row.get("evidence_title")),
            normalize(row.get("evidence_text")),
            normalize(row.get("evidence_date")),
            normalize(row.get("confidence_score"))
        ])

    ws4 = wb.create_sheet("citations_raw")
    ws4.append(["id", "claim_code", "evidence_item_id", "source_snapshot_id", "source_id", "citation_label", "citation_url", "citation_note"])
    for row in citation_rows:
        ws4.append([
            normalize(row.get("id")),
            normalize(row.get("claim_code")),
            normalize(row.get("evidence_item_id")),
            normalize(row.get("source_snapshot_id")),
            normalize(row.get("source_id")),
            normalize(row.get("citation_label")),
            normalize(row.get("citation_url")),
            normalize(row.get("citation_note"))
        ])

    ws5 = wb.create_sheet("snapshots_raw")
    ws5.append(["id", "company_id", "source_id", "snapshot_title", "snapshot_url", "raw_storage_path", "parsed_text_storage_path"])
    for row in snapshot_rows:
        ws5.append([
            normalize(row.get("id")),
            normalize(row.get("company_id")),
            normalize(row.get("source_id")),
            normalize(row.get("snapshot_title")),
            normalize(row.get("snapshot_url")),
            normalize(row.get("raw_storage_path")),
            normalize(row.get("parsed_text_storage_path"))
        ])

    ws6 = wb.create_sheet("metrics_raw")
    ws6.append(["id", "company_id", "metric_code", "metric_name", "metric_value", "metric_unit", "observation_date", "source_snapshot_id"])
    for row in metric_rows:
        ws6.append([
            normalize(row.get("id")),
            normalize(row.get("company_id")),
            normalize(row.get("metric_code")),
            normalize(row.get("metric_name")),
            normalize(row.get("metric_value")),
            normalize(row.get("metric_unit")),
            normalize(row.get("observation_date")),
            normalize(row.get("source_snapshot_id"))
        ])

    wb.save(path)

def make_pptx(path: Path, case: dict, sections: dict):
    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_bullets(slide, "Akkodis × Huawei Cloud 商业洞察", sections.get("executive_summary"))

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_bullets(slide, "行业层", sections.get("industry_view"))

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_bullets(slide, "AI技术栈层", sections.get("technology_stack_view"))

    for customer_name, items in (sections.get("customer_scene_view") or {}).items():
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        add_bullets(slide, f"{customer_name} 客户场景", items)

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_bullets(slide, "竞争格局层", sections.get("competitive_view"))

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_bullets(slide, "联合方案建议", sections.get("joint_solution_recommendations"))

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_bullets(slide, "风险与下一步", sections.get("risks_and_next_steps"))

    prs.save(path)

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--input-json-path", required=True)
    args = ap.parse_args()

    case = resolve_case_payload(args.input_json_path)

    evidence_rows, citation_rows, snapshot_rows, metric_rows = [], [], [], []
    collect_nested_rows(case, "recent_evidence_items", evidence_rows)
    collect_nested_rows(case, "recent_citations", citation_rows)
    collect_nested_rows(case, "recent_snapshots", snapshot_rows)
    collect_nested_rows(case, "recent_metric_observations", metric_rows)

    evidence_rows = evidence_rows[:500]
    citation_rows = citation_rows[:500]
    snapshot_rows = snapshot_rows[:500]
    metric_rows = metric_rows[:500]

    synth = build_synthesis(case, evidence_rows, citation_rows)
    sections = synth["sections"]

    case_code = case.get("case_code") or ((case.get("request") or {}).get("case_code")) or "account-strategy-case"
    out_dir = ROOT / "outputs" / "report_build" / case_code / ts_compact()
    out_dir.mkdir(parents=True, exist_ok=True)

    docx_path = out_dir / "account_strategy_report.docx"
    xlsx_path = out_dir / "account_strategy_workbook.xlsx"
    pptx_path = out_dir / "account_strategy_deck.pptx"
    json_path = out_dir / "account_strategy_bundle.json"

    make_docx(docx_path, case, sections, evidence_rows, citation_rows)
    make_xlsx(xlsx_path, case, evidence_rows, citation_rows, snapshot_rows, metric_rows)
    make_pptx(pptx_path, case, sections)

    out = {
        "ok": True,
        "skill_name": "report_build_skill.account_strategy_runner",
        "schema_version": "v2",
        "ts": utc_now(),
        "request": {
            "analysis_type": "account_strategy",
            "input_json_path": args.input_json_path,
            "case_code": case.get("case_code")
        },
        "conclusion": "rich account-strategy bundle generated with raw evidence/citation workbook and executive PPT.",
        "synthesis_mode": synth.get("mode"),
        "core_data": [
            {"field": "case_code", "value": case.get("case_code")},
            {"field": "evidence_count_total", "value": len(evidence_rows)},
            {"field": "citation_count_total", "value": len(citation_rows)},
            {"field": "snapshot_count_total", "value": len(snapshot_rows)},
            {"field": "metric_count_total", "value": len(metric_rows)}
        ],
        "artifact_plan": [
            {"artifact_type": "word", "path": str(docx_path.relative_to(ROOT))},
            {"artifact_type": "excel", "path": str(xlsx_path.relative_to(ROOT))},
            {"artifact_type": "ppt", "path": str(pptx_path.relative_to(ROOT))},
            {"artifact_type": "json", "path": str(json_path.relative_to(ROOT))}
        ]
    }

    write_json(json_path, out)
    print(json.dumps(out, ensure_ascii=False, indent=2))

if __name__ == "__main__":
    main()
