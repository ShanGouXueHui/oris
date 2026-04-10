#!/usr/bin/env python3
import argparse
import json
import re
from datetime import datetime, timezone
from pathlib import Path

from docx import Document
from docx.shared import Pt
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

import sys
ROOT = Path(__file__).resolve().parents[2]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from scripts.lib.oris_llm_client import call_oris_text

ROOT = Path(__file__).resolve().parents[2]
QUALITY_CFG_PATH = ROOT / "config" / "company_profile_quality.json"
FOCUS_PROFILE_PATH = ROOT / "config" / "company_focus_profiles.json"
METRIC_TAXONOMY_PATH = ROOT / "config" / "company_metric_taxonomy.json"
THEME_CFG_PATH = ROOT / "config" / "presentation_theme.json"
RULE_CFG_PATH = ROOT / "config" / "company_profile_rule_config.json"

def utc_now():
    return datetime.now(timezone.utc).isoformat()

def ts_compact():
    return datetime.now(timezone.utc).strftime("%Y%m%d_%H%M%S")

def load_json(path: Path):
    return json.loads(path.read_text(encoding="utf-8"))

def write_json(path: Path, data: dict):
    path.write_text(json.dumps(data, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")

def normalize(v):
    if v is None:
        return ""
    if isinstance(v, (dict, list)):
        return json.dumps(v, ensure_ascii=False)
    return str(v)

def safe_json_from_text(text: str):
    text = (text or "").strip()
    if not text:
        return None
    try:
        return json.loads(text)
    except Exception:
        pass
    import re
    m = re.search(r'(\{.*\})', text, flags=re.S)
    if m:
        try:
            return json.loads(m.group(1))
        except Exception:
            return None
    return None

def load_quality_cfg():
    try:
        return load_json(QUALITY_CFG_PATH)
    except Exception:
        return {}

def load_focus_profiles():
    try:
        return load_json(FOCUS_PROFILE_PATH)
    except Exception:
        return {"default_profile": "generic_company", "profiles": {}}

def pick_focus_profile(data: dict):
    req = data.get("request") or {}
    profile = (
        req.get("focus_profile")
        or data.get("focus_profile")
        or ((data.get("company_profile") or {}).get("focus_profile"))
        or ((data.get("db_backed_profile") or {}).get("focus_profile"))
    )
    if profile:
        return profile
    cfg = load_focus_profiles()
    return cfg.get("default_profile") or "generic_company"

def unified_profile(data: dict):
    profile = data.get("company_profile") or {}
    if profile:
        return profile
    profile = data.get("db_backed_profile") or {}
    if profile:
        return profile
    return {
        "company": {
            "company_name": (data.get("request") or {}).get("company_name")
        },
        "recent_snapshots": data.get("recent_snapshots") or [],
        "recent_evidence_items": data.get("recent_evidence_items") or [],
        "recent_metric_observations": data.get("recent_metric_observations") or []
    }

def pick_company_name(data: dict, profile: dict):
    company = (profile.get("company") or {})
    return (
        company.get("company_name")
        or (data.get("request") or {}).get("company_name")
        or "company"
    )

def flatten_text(v):
    if v is None:
        return ""
    if isinstance(v, str):
        return v.strip()
    if isinstance(v, list):
        return "\n".join([x for x in [flatten_text(i) for i in v] if x]).strip()
    if isinstance(v, dict):
        parts = []
        for k, val in v.items():
            t = flatten_text(val)
            if t:
                parts.append(f"{k}: {t}")
        return "\n".join(parts).strip()
    return str(v)

def score_text(text: str, focus_profile: str, cfg: dict):
    import re
    t = (text or "").strip()
    if not t:
        return -999
    tl = t.lower()
    score = 0

    if re.search(r"\d", t):
        score += 2
    if "%" in t:
        score += 2
    if re.search(r"(€|\$|¥|亿元|万台|million|billion|bn|mn)", t, flags=re.I):
        score += 2

    for kw in (cfg.get("finance_keywords") or []):
        if kw and kw.lower() in tl:
            score += 1

    prof_keywords = ((cfg.get("profile_keywords") or {}).get(focus_profile) or [])
    for kw in prof_keywords:
        if kw and kw.lower() in tl:
            score += 2

    for noise in (cfg.get("noise_patterns") or []):
        if noise and noise.lower() in tl:
            score -= 4

    if len(t) >= int(cfg.get("min_text_length", 35)):
        score += 1
    if len(t) > int(cfg.get("max_text_length", 320)):
        score -= 1

    return score

def top_evidence_items(profile: dict, focus_profile: str, cfg: dict, limit: int = 8):
    rows = []
    for row in profile.get("recent_evidence_items") or []:
        text = flatten_text(row.get("evidence_text") or "")
        title = flatten_text(row.get("evidence_title") or "")
        joined = f"{title}\n{text}".strip()
        if not joined:
            continue
        sc = score_text(joined, focus_profile, cfg)
        rows.append((sc, title, text))
    rows.sort(key=lambda x: (-x[0], -(len(x[2]) if x[2] else 0)))
    out = []
    for sc, title, text in rows[:limit]:
        body = text if len(text) <= 320 else text[:320] + "..."
        title = clean_evidence_label(title)
        body = clean_evidence_label(body)
        joined = f"{title}：{body}" if title else body
        if not is_user_facing_noise(joined, cfg):
            out.append(joined)
    return out

def top_metric_items(profile: dict, focus_profile: str, cfg: dict, limit: int = 8):
    rows = []
    blocklist = set(cfg.get("metric_blocklist") or [])
    for row in profile.get("recent_metric_observations") or []:
        code = normalize(row.get("metric_code"))
        if code in blocklist:
            continue
        name = normalize(row.get("metric_name"))
        value = normalize(row.get("metric_value"))
        unit = normalize(row.get("metric_unit"))
        obs = normalize(row.get("observation_date"))
        if not name and not code:
            continue
        text = f"{name or code}: {value}{unit} ({obs})".strip()
        sc = score_text(text, focus_profile, cfg) + 1
        rows.append((sc, text))
    rows.sort(key=lambda x: (-x[0], -len(x[1])))
    return [x[1] for x in rows[:limit]]

def profile_hint(focus_profile: str, cfg_profiles: dict):
    profiles = cfg_profiles.get("profiles") or {}
    prof = profiles.get(focus_profile) or {}
    label = prof.get("label") or focus_profile
    sections = prof.get("recommended_sections") or []
    req_metrics = prof.get("required_metric_codes") or []
    return {
        "profile_label": label,
        "recommended_sections": sections,
        "required_metric_codes": req_metrics
    }


def clean_evidence_label(text: str):
    import re
    t = normalize(text).strip()
    t = re.sub(r"\s*/\s*(derived|segment)\s*\d+\s*[:：]\s*", "：", t, flags=re.I)
    t = re.sub(r"\s+", " ", t).strip()
    return t

def is_internal_diagnostic_text(text: str):
    t = normalize(text).strip().lower()
    if not t:
        return True
    patterns = [
        "db-backed profile ready",
        "company row found in insight.company",
        "recent analysis_run count",
        "recent source_snapshot count",
        "recent evidence_item count",
        "recent metric_observation count",
        "company_profile_skill",
        "profile view:",
    ]
    return any(p in t for p in patterns)


def is_user_facing_noise(text: str, cfg: dict):
    t = normalize(text).strip()
    if not t:
        return True
    tl = t.lower()

    noise_patterns = [
        "investor alert options",
        "at least one of the checkboxes needs to be selected",
        "checkbox",
        "email address",
        "unsubscribe",
        "cookie",
        "cookies",
        "privacy",
        "newsletter",
        "accept all",
        "manage cookies",
        "opens in new window",
        "view all",
        "skip to content",
        "back to top",
        "news (includes earnings date announcements",
        "earnings date announcements",
        "upcoming conference appearances",
        "your information will be processed",
        "selecting a year value",
        "results公示",
        "结果公示",
        "会员大会",
        "章程（草案）",
        "筹备工作报告",
    ]
    for ptn in noise_patterns:
        if ptn in tl:
            return True

    for ptn in (cfg.get("noise_patterns") or []):
        if ptn and ptn.lower() in tl:
            return True

    return False

def filter_user_facing_items(items, cfg: dict, max_items: int = 8):
    out = []
    seen = set()
    for item in items or []:
        t = normalize(item).strip()
        if not t:
            continue
        if is_internal_diagnostic_text(t):
            continue
        if is_user_facing_noise(t, cfg):
            continue
        key = t.lower()
        if key in seen:
            continue
        seen.add(key)
        out.append(t)
        if len(out) >= max_items:
            break
    return out

def bucket_evidence_items(items, focus_profile: str):
    company, tech, industry, risk = [], [], [], []

    tech_kws = [
        "ai", "模型", "model", "agent", "api", "云", "cloud", "platform",
        "mb.os", "software-defined", "adas", "automated driving",
        "token", "推理", "训练", "open platform", "maas"
    ]
    industry_kws = [
        "market", "competition", "competitive", "industry", "ecosystem",
        "advertising", "games", "payment", "fintech", "cloud revenue",
        "市场", "竞争", "行业", "生态", "广告", "游戏", "支付"
    ]
    risk_kws = [
        "risk", "risks", "uncertainty", "challenge", "pressure", "regulation",
        "headwind", "decline", "slower", "volatile",
        "风险", "不确定", "挑战", "压力", "监管", "下滑", "波动"
    ]

    if focus_profile == "automotive_oem":
        tech_kws += ["electric", "vans", "xev", "range", "800v", "智驾", "电动车", "车型"]
    elif focus_profile == "internet_platform":
        industry_kws += ["mau", "ecosystem", "social", "ad", "广告", "用户", "社交"]
    elif focus_profile == "foundation_model_company":
        tech_kws += ["glm", "reasoning", "inference", "foundation model", "大模型", "推理"]

    for item in items or []:
        tl = item.lower()
        if any(k.lower() in tl for k in risk_kws):
            risk.append(item)
        elif any(k.lower() in tl for k in tech_kws):
            tech.append(item)
        elif any(k.lower() in tl for k in industry_kws):
            industry.append(item)
        else:
            company.append(item)

    return company, tech, industry, risk

def merge_unique(*groups, limit: int = 8):
    out = []
    seen = set()
    for group in groups:
        for item in group or []:
            t = normalize(item).strip()
            if not t:
                continue
            key = t.lower()
            if key in seen:
                continue
            seen.add(key)
            out.append(t)
            if len(out) >= limit:
                return out
    return out


def load_metric_taxonomy():
    try:
        return load_json(METRIC_TAXONOMY_PATH)
    except Exception:
        return {"version": 1, "profiles": {}}

def load_theme_cfg():
    try:
        return load_json(THEME_CFG_PATH)
    except Exception:
        return {
            "version": 1,
            "default_theme": "executive_blue",
            "themes": {
                "executive_blue": {
                    "palette": {
                        "bg": "#F7F9FC",
                        "title": "#0F2747",
                        "text": "#1F2937",
                        "muted": "#6B7280",
                        "accent": "#2563EB",
                        "accent_2": "#0EA5E9",
                        "risk": "#DC2626",
                        "success": "#059669",
                        "card_bg": "#FFFFFF",
                        "divider": "#D1D5DB"
                    }
                }
            }
        }


def load_rule_cfg():
    try:
        return load_json(RULE_CFG_PATH)
    except Exception:
        return {
            "version": 1,
            "default_theme_by_focus_profile": {"default": "executive_blue"},
            "user_facing_noise_patterns": [],
            "profile_negative_patterns": {"default": []},
            "profile_positive_metric_keywords": {"default": []},
            "polluted_external_segment_patterns": {"default": []}
        }

def rule_list(cfg: dict, key: str, focus_profile: str):
    section = cfg.get(key) or {}
    if isinstance(section, list):
        return section
    out = []
    out.extend(section.get("default") or [])
    out.extend(section.get(focus_profile) or [])
    seen = set()
    dedup = []
    for x in out:
        s = str(x).strip()
        if not s:
            continue
        k = s.lower()
        if k in seen:
            continue
        seen.add(k)
        dedup.append(s)
    return dedup

def hex_to_rgb(hex_str: str):
    s = (hex_str or "#000000").strip().lstrip("#")
    if len(s) != 6:
        s = "000000"
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))

def pick_theme_name(focus_profile: str, theme_cfg: dict):
    rule_cfg = load_rule_cfg()
    mapping = rule_cfg.get("default_theme_by_focus_profile") or {}
    return (
        mapping.get(focus_profile)
        or mapping.get("default")
        or theme_cfg.get("default_theme")
        or "executive_blue"
    )

def theme_palette(focus_profile: str):
    cfg = load_theme_cfg()
    theme_name = pick_theme_name(focus_profile, cfg)
    theme = (cfg.get("themes") or {}).get(theme_name) or {}
    palette = theme.get("palette") or {}
    return theme_name, palette

def safe_conclusion_line(data: dict):
    line = normalize(data.get("conclusion")).strip()
    if not line:
        return "公司洞察已生成。"
    if is_internal_diagnostic_text(line):
        return "已完成基于当前可得证据的公司洞察与结构化整理。"
    return line

def metric_taxonomy_for_profile(focus_profile: str):
    cfg = load_metric_taxonomy()
    profiles = cfg.get("profiles") or {}
    prof = profiles.get(focus_profile) or profiles.get("generic_company") or {}
    ordered = []
    for key in ["financial_core", "product_core", "operating_core", "model_core", "platform_core"]:
        ordered.extend(prof.get(key) or [])
    out = []
    seen = set()
    for x in ordered:
        if x not in seen:
            out.append(x)
            seen.add(x)
    return out

def metric_display_line(row: dict):
    name = normalize(row.get("metric_name") or row.get("metric_code"))
    value = normalize(row.get("metric_value"))
    unit = normalize(row.get("metric_unit"))
    obs = normalize(row.get("observation_date"))
    if not name:
        return ""
    tail = f"（{obs}）" if obs else ""
    return f"{name}: {value}{unit}{tail}".strip()


def normalize_metric_unit(unit: str):
    u = normalize(unit).strip()
    mapping = {
        "USD_billion": "亿美元",
        "million_accounts": "百万账户",
        "billion_users": "十亿用户",
        "percent": "%",
        "count": "个",
    }
    return mapping.get(u, u)

def metric_display_line_cn(row: dict):
    name = normalize(row.get("metric_name") or row.get("metric_code"))
    value = normalize(row.get("metric_value"))
    unit = normalize_metric_unit(row.get("metric_unit"))
    obs = normalize(row.get("observation_date"))
    if not name or value == "":
        return ""
    tail = f"（{obs}）" if obs else ""
    return f"{name}: {value}{unit}{tail}".strip()

def top_structured_metric_lines(profile: dict, focus_profile: str, limit: int = 8):
    preferred = metric_taxonomy_for_profile(focus_profile)
    rows = []
    block = {
        "official_source_snapshot_count",
        "extracted_evidence_segment_count",
    }
    for row in profile.get("recent_metric_observations") or []:
        code = normalize(row.get("metric_code"))
        if code in block:
            continue
        line = metric_display_line_cn(row)
        if not line:
            continue
        rows.append((code, line, row))

    positive_codes = {
        "revenue", "revenue_yoy", "alphabet_revenue", "alphabet_revenue_yoy",
        "google_cloud_revenue", "google_cloud_revenue_yoy",
        "google_services_revenue", "google_services_revenue_yoy",
        "gross_profit", "gross_profit_yoy",
        "operating_profit", "operating_profit_yoy",
        "net_profit", "free_cash_flow",
        "paid_subscriptions", "gemini_users", "advertising_revenue_share",
        "cloud_revenue", "cloud_revenue_yoy", "services_revenue", "services_revenue_yoy",
        "domestic_games_revenue", "domestic_games_revenue_yoy",
        "social_networks_revenue", "social_networks_revenue_yoy",
    }

    rank = {str(x).lower(): i for i, x in enumerate(preferred)}

    def sort_key(x):
        code, line, row = x
        c = code.lower()
        pos = 0 if c in positive_codes else 1
        return (pos, rank.get(c, 9999), -len(line))

    rows.sort(key=sort_key)

    out = []
    seen = set()
    for code, line, row in rows:
        k = line.lower()
        if k in seen:
            continue
        seen.add(k)
        out.append(line)
        if len(out) >= limit:
            break
    return out

def top_metric_items_by_taxonomy(profile: dict, focus_profile: str, cfg: dict, limit: int = 8):
    preferred = metric_taxonomy_for_profile(focus_profile)
    blocklist = set(cfg.get("metric_blocklist") or [])
    scored = []
    fallback = []

    for row in profile.get("recent_metric_observations") or []:
        code = normalize(row.get("metric_code"))
        if code in blocklist:
            continue
        text = metric_display_line(row)
        if not text:
            continue
        sc = score_text(text, focus_profile, cfg) + 1
        payload = (sc, text, code, row)
        if metric_matches_taxonomy(row, preferred):
            scored.append(payload)
        else:
            fallback.append(payload)

    code_rank = {code.lower(): idx for idx, code in enumerate(preferred)}
    scored.sort(key=lambda x: (code_rank.get(normalize(x[2]).lower(), 9999), -x[0], -len(x[1])))
    fallback.sort(key=lambda x: (-x[0], -len(x[1])))

    out = []
    seen = set()
    for group in [scored, fallback]:
        for sc, text, code, row in group:
            key = text.lower()
            if key in seen:
                continue
            seen.add(key)
            out.append(text)
            if len(out) >= limit:
                return out
    return out

def top_metric_cards(profile: dict, focus_profile: str, cfg: dict, limit: int = 4):
    preferred = metric_taxonomy_for_profile(focus_profile)
    rows = []
    for row in profile.get("recent_metric_observations") or []:
        if not metric_matches_taxonomy(row, preferred):
            continue
        code = normalize(row.get("metric_code"))
        title = normalize(row.get("metric_name") or code)
        value = normalize(row.get("metric_value"))
        unit = normalize(row.get("metric_unit"))
        obs = normalize(row.get("observation_date"))
        if not title or not value:
            continue
        rows.append({
            "metric_code": code,
            "title": title,
            "value": f"{value}{unit}".strip(),
            "subtitle": obs or ""
        })
    rank = {code.lower(): idx for idx, code in enumerate(preferred)}
    rows.sort(key=lambda x: rank.get(normalize(x["metric_code"]).lower(), 9999))
    return rows[:limit]

def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=None, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if hasattr(p, "clear"):
        p.clear()
    p.alignment = align
    run = p.add_run()
    safe_text = str(text) if text is not None else ""
    run.text = safe_text if safe_text.strip() else " "
    run.font.size = PptPt(font_size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    return box

def add_rect(slide, left, top, width, height, fill_rgb, line_rgb=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.color.rgb = line_rgb or fill_rgb
    return shape

def apply_slide_bg(slide, bg_rgb):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = bg_rgb

def make_pptx_cover(prs, company_name: str, focus_profile: str, palette: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_bg(slide, hex_to_rgb(palette.get("bg", "#F7F9FC")))
    add_textbox(slide, Inches(0.7), Inches(1.0), Inches(11.0), Inches(0.8),
                f"{company_name} 公司洞察", font_size=28, bold=True, color=hex_to_rgb(palette.get("title", "#0F2747")))
    add_textbox(slide, Inches(0.75), Inches(1.9), Inches(10.5), Inches(0.5),
                f"焦点画像：{focus_profile}", font_size=14, bold=False, color=hex_to_rgb(palette.get("muted", "#6B7280")))
    add_rect(slide, Inches(0.75), Inches(2.5), Inches(3.2), Inches(0.08), hex_to_rgb(palette.get("accent", "#2563EB")))
    return slide

def make_pptx_exec_summary(prs, sections: dict, palette: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_bg(slide, hex_to_rgb(palette.get("bg", "#F7F9FC")))
    add_textbox(slide, Inches(0.6), Inches(0.4), Inches(6), Inches(0.5),
                "执行摘要", font_size=24, bold=True, color=hex_to_rgb(palette.get("title", "#0F2747")))
    items = (sections.get("executive_summary") or [])[:3]
    card_w = Inches(3.7)
    xs = [Inches(0.6), Inches(4.3), Inches(8.0)]
    for i in range(3):
        add_rect(slide, xs[i], Inches(1.2), card_w, Inches(3.0),
                 hex_to_rgb(palette.get("card_bg", "#FFFFFF")),
                 hex_to_rgb(palette.get("divider", "#D1D5DB")))
        title = ["核心判断", "关键证据", "持续跟踪"][i]
        body = items[i] if i < len(items) else "N/A"
        add_textbox(slide, xs[i] + Inches(0.2), Inches(1.45), Inches(3.1), Inches(0.35),
                    title, font_size=13, bold=True, color=hex_to_rgb(palette.get("accent", "#2563EB")))
        add_textbox(slide, xs[i] + Inches(0.2), Inches(1.85), Inches(3.15), Inches(2.1),
                    body, font_size=15, bold=False, color=hex_to_rgb(palette.get("text", "#1F2937")))
    return slide

def make_pptx_kpi_dashboard(prs, metric_cards: list, palette: dict):
    if not metric_cards:
        return
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_bg(slide, hex_to_rgb(palette.get("bg", "#F7F9FC")))
    add_textbox(slide, Inches(0.6), Inches(0.4), Inches(6), Inches(0.5),
                "关键指标仪表板", font_size=24, bold=True, color=hex_to_rgb(palette.get("title", "#0F2747")))
    positions = [
        (Inches(0.7), Inches(1.3)),
        (Inches(3.5), Inches(1.3)),
        (Inches(6.3), Inches(1.3)),
        (Inches(9.1), Inches(1.3)),
    ]
    for idx, card in enumerate(metric_cards[:4]):
        left, top = positions[idx]
        add_rect(slide, left, top, Inches(2.3), Inches(2.2),
                 hex_to_rgb(palette.get("card_bg", "#FFFFFF")),
                 hex_to_rgb(palette.get("divider", "#D1D5DB")))
        add_textbox(slide, left + Inches(0.15), top + Inches(0.18), Inches(2.0), Inches(0.35),
                    card.get("title") or "指标", font_size=12, bold=True, color=hex_to_rgb(palette.get("muted", "#6B7280")))
        add_textbox(slide, left + Inches(0.15), top + Inches(0.75), Inches(2.0), Inches(0.55),
                    card.get("value") or "-", font_size=22, bold=True, color=hex_to_rgb(palette.get("accent", "#2563EB")))
        add_textbox(slide, left + Inches(0.15), top + Inches(1.55), Inches(2.0), Inches(0.3),
                    card.get("subtitle") or "", font_size=10, bold=False, color=hex_to_rgb(palette.get("muted", "#6B7280")))
    return slide

def make_pptx_two_column(prs, title: str, items: list, palette: dict):
    if not items:
        return
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_bg(slide, hex_to_rgb(palette.get("bg", "#F7F9FC")))
    add_textbox(slide, Inches(0.6), Inches(0.4), Inches(8), Inches(0.5),
                title, font_size=24, bold=True, color=hex_to_rgb(palette.get("title", "#0F2747")))
    cols = [items[::2], items[1::2]]
    for ci, col in enumerate(cols):
        left = Inches(0.7) if ci == 0 else Inches(6.4)
        add_rect(slide, left, Inches(1.2), Inches(5.1), Inches(5.5),
                 hex_to_rgb(palette.get("card_bg", "#FFFFFF")),
                 hex_to_rgb(palette.get("divider", "#D1D5DB")))
        y = 1.45
        for item in col[:6]:
            add_textbox(slide, left + Inches(0.18), Inches(y), Inches(4.7), Inches(0.7),
                        f"• {item}", font_size=13, bold=False, color=hex_to_rgb(palette.get("text", "#1F2937")))
            y += 0.82
    return slide

def make_pptx_risk_matrix(prs, items: list, palette: dict):
    if not items:
        return
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_bg(slide, hex_to_rgb(palette.get("bg", "#F7F9FC")))
    add_textbox(slide, Inches(0.6), Inches(0.4), Inches(6), Inches(0.5),
                "风险矩阵", font_size=24, bold=True, color=hex_to_rgb(palette.get("title", "#0F2747")))
    for idx, item in enumerate(items[:4]):
        left = Inches(0.8 + (idx % 2) * 5.6)
        top = Inches(1.3 + (idx // 2) * 2.4)
        add_rect(slide, left, top, Inches(4.8), Inches(1.8),
                 hex_to_rgb(palette.get("card_bg", "#FFFFFF")),
                 hex_to_rgb(palette.get("risk", "#DC2626")))
        add_textbox(slide, left + Inches(0.18), top + Inches(0.18), Inches(4.3), Inches(1.2),
                    item, font_size=13, bold=False, color=hex_to_rgb(palette.get("text", "#1F2937")))
    return slide

def make_pptx_tracking(prs, items: list, palette: dict):
    if not items:
        return
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_bg(slide, hex_to_rgb(palette.get("bg", "#F7F9FC")))
    add_textbox(slide, Inches(0.6), Inches(0.4), Inches(6), Inches(0.5),
                "持续跟踪指标", font_size=24, bold=True, color=hex_to_rgb(palette.get("title", "#0F2747")))
    y = 1.2
    for item in items[:8]:
        add_rect(slide, Inches(0.8), Inches(y), Inches(11.2), Inches(0.55),
                 hex_to_rgb(palette.get("card_bg", "#FFFFFF")),
                 hex_to_rgb(palette.get("divider", "#D1D5DB")))
        add_textbox(slide, Inches(1.0), Inches(y + 0.1), Inches(10.8), Inches(0.3),
                    item, font_size=12, bold=False, color=hex_to_rgb(palette.get("text", "#1F2937")))
        y += 0.68
    return slide


def is_engineering_next_step(text: str):
    t = normalize(text).strip().lower()
    if not t:
        return True
    patterns = [
        "increase official source coverage",
        "replace bootstrap evidence",
        "let report_build_skill consume",
        "db-backed",
        "analysis_run",
        "source_snapshot",
        "evidence_item",
        "metric rows",
        "persisted snapshots",
        "downstream report assembly",
        "request-only payload",
        "bootstrap evidence",
        "official source coverage"
    ]
    return any(p in t for p in patterns)

def taxonomy_aliases():
    return {
        "revenue": ["revenue", "收入", "营收", "sales"],
        "gross_profit": ["gross profit", "毛利"],
        "operating_profit": ["operating profit", "ebit", "营业利润", "经营利润"],
        "net_profit": ["net profit", "profit attributable", "净利润", "净利"],
        "free_cash_flow": ["free cash flow", "自由现金流"],
        "r_and_d_expense": ["r&d", "research and development", "研发", "研发费用"],
        "cloud_revenue": ["cloud revenue", "云收入", "云业务收入"],
        "advertising_revenue": ["advertising revenue", "广告收入"],
        "api_revenue": ["api revenue", "api收入"],
        "monthly_tokens": ["tokens", "monthly tokens", "token"],
        "enterprise_customer_count": ["enterprise customer", "客户数量", "企业客户"],
        "mau": ["mau", "monthly active users", "月活"],
        "dau": ["dau", "daily active users", "日活"],
        "nps": ["nps", "net promoter score"],
        "conversion_rate": ["conversion", "转化率"],
        "vehicle_sales_total": ["vehicle sales", "deliveries", "销量", "交付"],
        "ev_sales": ["ev sales", "新能源销量", "电动车销量"],
        "range_km": ["range", "续航"],
        "acceleration_0_100_s": ["0-100", "0 to 100", "百公里加速", "百米加速"],
        "adas_takeover_per_1000km": ["takeover", "接管", "接手次数"],
    }

def metric_matches_taxonomy(row: dict, preferred_codes: list[str]):
    code = normalize(row.get("metric_code")).strip().lower()
    name = normalize(row.get("metric_name")).strip().lower()
    text = f"{code} {name}".strip()

    if code in {x.lower() for x in preferred_codes}:
        return True

    alias_map = taxonomy_aliases()
    for pref in preferred_codes:
        aliases = alias_map.get(pref, []) + [pref]
        for alias in aliases:
            if alias and alias.lower() in text:
                return True
    return False



def split_numeric_segments(text: str):
    t = normalize(text).strip()
    if not t:
        return []
    raw_parts = re.split(r"[\n\r]+|[•●▪◦]|\s{2,}|(?<=；)\s+|(?<=;)\s+|(?<=\.)\s+(?=[A-Z0-9])", t)
    out = []
    for part in raw_parts:
        seg = normalize(part).strip(" -•●▪◦")
        seg = re.sub(r"\s+", " ", seg).strip()
        if len(seg) < 18:
            continue
        out.append(seg)
    return out



def derived_kpis_from_evidence(profile: dict, focus_profile: str, cfg: dict, limit: int = 8):
    rows = []
    negs = [x.lower() for x in profile_negative_patterns(focus_profile)]

    for row in profile.get("recent_evidence_items") or []:
        title = clean_evidence_label(flatten_text(row.get("evidence_title") or ""))
        body = flatten_text(row.get("evidence_text") or "")
        if not body:
            continue

        for seg in split_numeric_segments(body):
            joined = f"{title}：{seg}" if title else seg
            jl = joined.lower()

            if is_internal_diagnostic_text(joined):
                continue
            if is_user_facing_noise(joined, cfg):
                continue
            if any(x in jl for x in negs):
                continue
            if looks_like_title_echo(seg, title):
                continue
            if is_polluted_external_segment(seg, focus_profile):
                continue
            if not is_metric_like_segment(seg, focus_profile):
                continue

            metric_line = compress_metric_segment(seg, title, focus_profile)
            if not metric_line:
                continue
            if is_user_facing_noise(metric_line, cfg):
                continue
            if looks_like_title_echo(metric_line, title):
                continue
            if is_polluted_external_segment(metric_line, focus_profile):
                continue

            sc = score_text(metric_line, focus_profile, cfg) + 3
            if re.search(r"\d", metric_line):
                sc += 1
            if re.search(r"(€|\$|¥|rmb|usd|亿元|万元|million|billion|bn|mn|mau|dau|tokens|销量|交付|用户|客户|revenue|profit|cash flow)", metric_line, flags=re.I):
                sc += 2
            rows.append((sc, metric_line))

    rows.sort(key=lambda x: (-x[0], -len(x[1])))
    out = []
    seen = set()
    for sc, seg in rows:
        key = seg.lower()
        if key in seen:
            continue
        seen.add(key)
        out.append(seg)
        if len(out) >= limit:
            break
    return out

def metric_cards_from_lines(lines, limit: int = 4):
    cards = []
    for line in lines or []:
        t = normalize(line).strip()
        if not t:
            continue

        t_clean = re.sub(r"（[^）]{0,40}）$", "", t).strip()
        title = "关键指标"
        value = t_clean
        subtitle = ""

        m1 = re.search(r"^(.{4,40}?)(?:收入|营收|利润|毛利|现金流|用户|客户|销量|交付|tokens?|Revenue|profit|users|customers)", t_clean, flags=re.I)
        if m1:
            title = m1.group(0)[:28]

        m2 = re.search(r"((?:RMB|USD|US\$|\$|€)?\s?[\d\.,]+(?:\s?(?:billion|million|bn|mn|%))?)", t_clean, flags=re.I)
        if m2:
            value = m2.group(1).strip()
            subtitle = t_clean[:56]
        else:
            value = t_clean[:42]

        cards.append({
            "metric_code": "",
            "title": title,
            "value": value,
            "subtitle": subtitle[:48]
        })
        if len(cards) >= limit:
            break
    return cards

def source_label_short(title: str):
    t = clean_evidence_label(title)
    t = re.sub(r"\s*/\s*(derived|segment)\s*\d+\s*$", "", t, flags=re.I)
    t = re.sub(r"\s+", " ", t).strip()
    return t

def profile_negative_patterns(focus_profile: str):
    rule_cfg = load_rule_cfg()
    return [x.lower() for x in rule_list(rule_cfg, "profile_negative_patterns", focus_profile)]


def is_metric_like_segment(seg: str, focus_profile: str):
    t = normalize(seg).strip()
    if not t:
        return False
    tl = t.lower()

    pos = profile_positive_metric_keywords(focus_profile)
    has_pos_kw = any(x in tl for x in pos)
    has_number = bool(re.search(r"\d", t))
    has_metric_unit = bool(re.search(r"(€|\$|¥|rmb|usd|亿元|万元|million|billion|bn|mn|%|mau|dau|tokens?)", t, flags=re.I))

    if focus_profile == "internet_platform":
        return (has_number and (has_pos_kw or has_metric_unit))

    if focus_profile == "foundation_model_company":
        bench_kw = any(x in tl for x in ["swe-bench", "terminal bench", "sota", "50 步", "50-step", "50 step", "dozens of", "benchmark"])
        return (has_number and (has_pos_kw or has_metric_unit or bench_kw)) or bench_kw

    if focus_profile == "automotive_oem":
        return has_number and (has_pos_kw or has_metric_unit)

    return has_number and (has_pos_kw or has_metric_unit)


def compress_metric_segment(seg: str, title: str, focus_profile: str):
    t = normalize(seg).strip()
    if not t:
        return ""
    src = source_label_short(title)

    if src and t.lower().startswith(src.lower()):
        t = t[len(src):].lstrip(" ：:-")

    t = re.sub(r"\s+", " ", t).strip()
    t = re.sub(r"\s*[\(\[]?(due to|reflecting|including|which was mainly due to|primarily driven by).*$", "", t, flags=re.I)
    t = re.sub(r"\s*your information will be processed.*$", "", t, flags=re.I)
    t = re.sub(r"\s*news \(includes earnings date announcements.*$", "", t, flags=re.I)

    parts = re.split(r"(?<=[\.\!\?。；;])\s+|,\s+(?=[A-Z0-9])|，\s*", t)
    kept = []
    for part in parts:
        s = normalize(part).strip(" -•●▪◦")
        if not s:
            continue
        kept.append(s)
        joined = "；".join(kept)
        if len(joined) >= 110:
            break

    out = "；".join(kept).strip() if kept else t
    out = re.sub(r"\s+", " ", out).strip("；;，, ")

    out = re.sub(r"^\d{1,2}\s*月\s*\d{1,2}\s*日[；;，,\s]*", "", out)
    out = re.sub(r"^20\d{2}年\d{1,2}月\d{1,2}日[；;，,\s]*", "", out)

    if len(out) > 135:
        out = out[:135].rstrip() + "..."
    if not out:
        return ""

    if src:
        return f"{out}（{src}）"
    return out

def profile_positive_metric_keywords(focus_profile: str):
    rule_cfg = load_rule_cfg()
    vals = rule_list(rule_cfg, "profile_positive_metric_keywords", focus_profile)
    return [x.lower() for x in vals]

def looks_like_title_echo(seg: str, title: str):
    s = normalize(seg).strip().lower()
    t = source_label_short(title).strip().lower()
    if not s or not t:
        return False
    if s == t:
        return True
    if s.startswith(t):
        return True
    if len(s) <= max(40, len(t) + 20) and t in s:
        return True
    return False

def is_polluted_external_segment(seg: str, focus_profile: str):
    s = normalize(seg).strip().lower()
    if not s:
        return True
    rule_cfg = load_rule_cfg()
    patterns = [x.lower() for x in rule_list(rule_cfg, "polluted_external_segment_patterns", focus_profile)]
    return any(p in s for p in patterns)

def _fmt_metric_display_number(x: float):
    try:
        x = float(x)
    except Exception:
        return str(x)
    if abs(x - round(x)) < 1e-9:
        return str(int(round(x)))
    return f"{x:.1f}".rstrip("0").rstrip(".")

def canonicalize_metric_line(line: str):
    t = normalize(line).strip()
    if not t:
        return ""

    m = re.match(r'^(.*?:)\s*([0-9]+(?:\.[0-9]+)?)([A-Za-z_]+)(（[^）]+）)$', t)
    if not m:
        return t

    head, num_s, unit, tail = m.groups()
    try:
        num = float(num_s)
    except Exception:
        return t

    head = head.strip()

    if unit == "percent":
        return f"{head} {_fmt_metric_display_number(num)}%{tail}"
    if unit == "USD_billion":
        return f"{head} {_fmt_metric_display_number(num)}亿美元{tail}"
    if unit == "USD_million":
        return f"{head} {_fmt_metric_display_number(num / 100.0)}亿美元{tail}"
    if unit == "million_accounts":
        return f"{head} {_fmt_metric_display_number(num)}百万账户{tail}"
    if unit == "billion_users":
        return f"{head} {_fmt_metric_display_number(num)}十亿用户{tail}"

    return t

def build_synthesis(data: dict):
    profile = unified_profile(data)
    focus_profile = pick_focus_profile(data)
    cfg = load_quality_cfg()
    cfg_profiles = load_focus_profiles()
    hint = profile_hint(focus_profile, cfg_profiles)

    evidence_items = profile.get("recent_evidence_items") or []
    metric_items = profile.get("recent_metric_observations") or []

    prompt = {
        "task": "生成高质量公司洞察内容",
        "requirements": {
            "style": "专业、深入、分层、偏证据优先",
            "must_cover": [
                "公司规模与覆盖区域",
                "主营业务与品牌结构",
                "核心产品与核心竞争力",
                "AI/智能化技术栈与产品方向",
                "行业位置与竞争态势",
                "风险与持续跟踪指标"
            ],
            "must_distinguish": ["事实", "推断", "建议", "风险"],
            "chat_first": True
        },
        "focus_profile": focus_profile,
        "profile_hint": hint,
        "input": {
            "request": data.get("request"),
            "core_data": data.get("core_data"),
            "facts": data.get("facts"),
            "inferences": data.get("inferences"),
            "risks": data.get("risks"),
            "next_steps": data.get("next_steps"),
            "recent_snapshots": profile.get("recent_snapshots"),
            "recent_evidence_items": evidence_items[:20],
            "recent_metric_observations": metric_items[:20]
        },
        "output_schema": {
            "executive_summary": ["..."],
            "company_view": ["..."],
            "technology_view": ["..."],
            "industry_view": ["..."],
            "risk_view": ["..."],
            "tracking_kpis": ["..."]
        }
    }

    try:
        resp = call_oris_text(
            "请基于下面JSON生成严格JSON，不要输出解释文字：\n" + json.dumps(prompt, ensure_ascii=False),
            role="free_fallback",
            timeout_seconds=180
        )
        if resp.get("ok"):
            parsed = safe_json_from_text(resp.get("text", ""))
            if parsed and isinstance(parsed, dict):
                return {"mode": "llm", "sections": parsed, "focus_profile": focus_profile}
    except Exception:
        pass

    evidence_top = top_evidence_items(profile, focus_profile, cfg, limit=int(cfg.get("max_evidence_items", 8)))
    metrics_top = top_structured_metric_lines(profile, focus_profile, limit=int(cfg.get("max_metric_items", 8)))
    metrics_top = [canonicalize_metric_line(x) for x in metrics_top if normalize(x).strip()]
    derived_kpis = derived_kpis_from_evidence(profile, focus_profile, cfg, limit=int(cfg.get("max_metric_items", 8)))
    if not metrics_top:
        metrics_top = top_structured_metric_lines(profile, focus_profile, limit=int(cfg.get("max_metric_items", 8)))
    if not metrics_top:
        metrics_top = derived_kpis

    clean_facts = filter_user_facing_items(data.get("facts") or [], cfg, max_items=8)
    clean_inferences = filter_user_facing_items(data.get("inferences") or [], cfg, max_items=8)
    clean_risks = filter_user_facing_items(data.get("risks") or [], cfg, max_items=8)
    clean_next_steps = [
        x for x in filter_user_facing_items(data.get("next_steps") or [], cfg, max_items=8)
        if not is_engineering_next_step(x)
    ]

    ev_company, ev_tech, ev_industry, ev_risk = bucket_evidence_items(evidence_top, focus_profile)

    summary = [
        safe_conclusion_line(data),
        f"当前焦点画像：{hint['profile_label']}。",
        "建议以业务结构、核心产品/技术、竞争位置、风险与跟踪指标五层持续跟踪。"
    ]
    if metrics_top:
        summary.append("优先关注指标：" + "；".join(metrics_top[:3]))
    elif derived_kpis:
        summary.append("当前可提炼量化片段：" + "；".join(derived_kpis[:2]))

    sections = {
        "executive_summary": merge_unique(summary, limit=4),
        "company_view": merge_unique(clean_facts, ev_company, limit=8),
        "technology_view": merge_unique(clean_inferences, ev_tech, limit=8),
        "industry_view": merge_unique(
            ["公司洞察不应停留在官网口径，需叠加行业、竞争、客户与资本市场视角形成完整判断。"],
            ev_industry,
            [f"建议重点展开的章节：{', '.join(hint['recommended_sections']) if hint['recommended_sections'] else 'company_positioning, business_model, risks, tracking_metrics'}。"],
            limit=8
        ),
        "risk_view": merge_unique(clean_risks, ev_risk, limit=8),
        "tracking_kpis": merge_unique(metrics_top, clean_next_steps, limit=8)
    }
    return {"mode": "deterministic_fallback", "sections": sections, "focus_profile": focus_profile}

def add_bullets(slide, title, items):
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    items = items or ["N/A"]
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = str(item)

def make_docx(path: Path, data: dict, sections: dict):
    doc = Document()
    profile = unified_profile(data)
    company_name = pick_company_name(data, profile)
    title = doc.add_heading(f"{company_name} 公司洞察报告", 0)
    title.runs[0].font.size = Pt(20)

    for name, title_text in [
        ("executive_summary", "执行摘要"),
        ("company_view", "公司层"),
        ("technology_view", "技术层"),
        ("industry_view", "行业层"),
        ("risk_view", "风险层"),
        ("tracking_kpis", "持续跟踪指标")
    ]:
        doc.add_heading(title_text, level=1)
        for item in sections.get(name) or []:
            doc.add_paragraph(str(item), style="List Bullet")

    doc.add_heading("原始证据", level=1)
    for row in (profile.get("recent_evidence_items") or [])[:30]:
        doc.add_paragraph(
            f"{normalize(row.get('evidence_title'))}\n{normalize(row.get('evidence_text'))}",
            style="List Bullet"
        )

    doc.save(path)

def make_xlsx(path: Path, data: dict):
    wb = Workbook()

    ws = wb.active
    ws.title = "core_data"
    ws.append(["field", "value"])
    for row in data.get("core_data") or []:
        ws.append([normalize(row.get("field")), normalize(row.get("value"))])

    ws2 = wb.create_sheet("sources")
    ws2.append(["source_name", "source_type", "url", "official_flag"])
    for row in data.get("sources") or []:
        ws2.append([
            normalize(row.get("source_name")),
            normalize(row.get("source_type")),
            normalize(row.get("url")),
            normalize(row.get("official_flag"))
        ])

    profile = unified_profile(data)

    ws3 = wb.create_sheet("evidence_raw")
    ws3.append(["id", "source_snapshot_id", "evidence_type", "evidence_title", "evidence_text", "evidence_date", "confidence_score"])
    for row in profile.get("recent_evidence_items") or []:
        ws3.append([
            normalize(row.get("id")),
            normalize(row.get("source_snapshot_id")),
            normalize(row.get("evidence_type")),
            normalize(row.get("evidence_title")),
            normalize(row.get("evidence_text")),
            normalize(row.get("evidence_date")),
            normalize(row.get("confidence_score"))
        ])

    ws4 = wb.create_sheet("metrics_raw")
    ws4.append(["id", "metric_code", "metric_name", "metric_value", "metric_unit", "observation_date", "source_snapshot_id"])
    for row in profile.get("recent_metric_observations") or []:
        ws4.append([
            normalize(row.get("id")),
            normalize(row.get("metric_code")),
            normalize(row.get("metric_name")),
            normalize(row.get("metric_value")),
            normalize(row.get("metric_unit")),
            normalize(row.get("observation_date")),
            normalize(row.get("source_snapshot_id"))
        ])

    wb.save(path)


def make_pptx(path: Path, sections: dict, company_name: str, focus_profile: str, metric_cards: list | None = None):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    theme_name, palette = theme_palette(focus_profile)

    make_pptx_cover(prs, company_name, focus_profile, palette)
    make_pptx_exec_summary(prs, sections, palette)
    make_pptx_kpi_dashboard(prs, metric_cards or [], palette)
    make_pptx_two_column(prs, "公司与商业模式", sections.get("company_view") or [], palette)
    make_pptx_two_column(prs, "技术与产品能力", sections.get("technology_view") or [], palette)
    make_pptx_two_column(prs, "行业与竞争", sections.get("industry_view") or [], palette)
    make_pptx_risk_matrix(prs, sections.get("risk_view") or [], palette)
    make_pptx_tracking(prs, sections.get("tracking_kpis") or [], palette)

    prs.save(path)


def metric_should_hide(row: dict):
    code = normalize(row.get("metric_code")).lower()
    unit = normalize(row.get("metric_unit"))
    try:
        value = float(row.get("metric_value"))
    except Exception:
        value = None

    if value is None:
        return True

    if code in {"official_source_snapshot_count", "extracted_evidence_segment_count"}:
        return True

    if code == "google_cloud_revenue" and unit == "USD_million" and value > 50000:
        return True

    if code == "google_services_revenue" and unit == "USD_million" and value > 500000:
        return True

    if code == "advertising_revenue_share" and not (0 <= value <= 100):
        return True

    if code.endswith("_yoy") and unit == "percent" and not (-100 <= value <= 1000):
        return True

    if code == "gemini_users" and not (0 < value <= 10):
        return True

    if code == "paid_subscriptions" and not (0 < value <= 5000):
        return True

    return False


def metric_priority_for_profile_v2(focus_profile: str):
    if focus_profile == "internet_platform":
        return [
            "google_services_revenue",
            "google_services_revenue_yoy",
            "google_cloud_revenue",
            "google_cloud_revenue_yoy",
            "google_search_and_other_revenue",
            "advertising_revenue_share",
            "paid_subscriptions",
            "gemini_users",
            "revenue",
            "revenue_yoy",
            "gross_profit",
            "gross_profit_yoy",
            "operating_profit",
            "operating_profit_yoy",
            "net_profit",
            "free_cash_flow",
        ]
    if focus_profile == "foundation_model_company":
        return [
            "api_revenue",
            "enterprise_customer_count",
            "monthly_tokens",
            "benchmark_score",
            "agent_task_steps",
            "gemini_users",
            "revenue",
            "revenue_yoy",
        ]
    if focus_profile == "automotive_oem":
        return [
            "vehicle_sales_total",
            "ev_sales",
            "revenue",
            "revenue_yoy",
            "gross_profit",
            "operating_profit",
            "net_profit",
            "free_cash_flow",
        ]
    return [
        "revenue", "revenue_yoy", "gross_profit", "operating_profit",
        "net_profit", "free_cash_flow"
    ]


def metric_rank_score_v2(row: dict, focus_profile: str):
    code = normalize(row.get("metric_code")).lower()
    unit = normalize(row.get("metric_unit"))
    try:
        value = float(row.get("metric_value"))
    except Exception:
        value = 0.0

    score = 0
    if unit in {"USD_billion", "RMB_billion", "percent", "million_accounts", "billion_users"}:
        score += 20
    elif unit in {"USD_million", "RMB_million"}:
        score += 10

    if code.endswith("_yoy") and unit == "percent":
        score += 6

    if code == "google_cloud_revenue" and unit == "USD_billion" and 0 < value < 500:
        score += 12
    if code == "google_services_revenue" and unit == "USD_billion" and 0 < value < 500:
        score += 12
    if code == "google_search_and_other_revenue" and unit == "USD_million" and 0 < value < 500000:
        score += 8

    if focus_profile == "internet_platform":
        if code in {"google_services_revenue", "google_cloud_revenue", "advertising_revenue_share", "paid_subscriptions", "gemini_users"}:
            score += 5

    return score


def top_structured_metric_lines(profile: dict, focus_profile: str, limit: int = 8):
    priority = metric_priority_for_profile_v2(focus_profile)
    rank = {x.lower(): i for i, x in enumerate(priority)}

    best_by_code = {}
    for row in profile.get("recent_metric_observations") or []:
        code = normalize(row.get("metric_code")).lower()
        if not code:
            continue
        if metric_should_hide(row):
            continue

        line = metric_display_line_cn(row)
        if not line:
            continue

        candidate = (
            rank.get(code, 9999),
            -metric_rank_score_v2(row, focus_profile),
            line,
            normalize(row.get("observation_date")),
            normalize(row.get("source_snapshot_id")),
            normalize(row.get("evidence_item_id")),
        )
        old = best_by_code.get(code)
        if old is None or candidate < old:
            best_by_code[code] = candidate

    rows = sorted(best_by_code.values(), key=lambda x: (x[0], x[1], x[2]))
    out = []
    for row in rows:
        line = row[2]
        out.append(line)
        if len(out) >= limit:
            break
    return out


def top_metric_cards(profile: dict, focus_profile: str, cfg: dict, limit: int = 4):
    lines = top_structured_metric_lines(profile, focus_profile, limit=limit)
    cards = []
    for line in lines:
        title = "关键指标"
        value = line
        subtitle = ""
        if "（" in line and line.endswith("）"):
            left, right = line.rsplit("（", 1)
            subtitle = right[:-1]
        else:
            left = line

        if ":" in left:
            title, value = left.split(":", 1)
            title = title.strip()[:28]
            value = value.strip()[:34]
        else:
            value = left.strip()[:34]

        cards.append({
            "metric_code": "",
            "title": title,
            "value": value,
            "subtitle": subtitle[:24],
        })
    return cards


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--input-json-path", required=True)
    args = ap.parse_args()

    data = load_json(Path(args.input_json_path))
    profile = unified_profile(data)
    company_name = pick_company_name(data, profile)

    synth = build_synthesis(data)
    sections = synth["sections"]
    focus_profile = synth.get("focus_profile") or pick_focus_profile(data)
    metric_cards = top_metric_cards(profile, focus_profile, load_quality_cfg(), limit=4)
    if not metric_cards:
        metric_cards = metric_cards_from_lines((sections.get("tracking_kpis") or [])[:4], limit=4)

    slug = company_name.lower().replace(" ", "-")
    out_dir = ROOT / "outputs" / "report_build" / f"company-profile-{slug}" / ts_compact()
    out_dir.mkdir(parents=True, exist_ok=True)

    docx_path = out_dir / "company_profile_report.docx"
    xlsx_path = out_dir / "company_profile_workbook.xlsx"
    pptx_path = out_dir / "company_profile_deck.pptx"
    json_path = out_dir / "company_profile_bundle.json"

    make_docx(docx_path, data, sections)
    make_xlsx(xlsx_path, data)
    make_pptx(pptx_path, sections, company_name, focus_profile, metric_cards)

    bundle = {
        "ok": True,
        "skill_name": "report_build_skill.company_profile_bundle_runner",
        "schema_version": "v3",
        "ts": utc_now(),
        "request": {
            "analysis_type": "company_profile",
            "input_json_path": args.input_json_path,
            "company_name": company_name,
            "focus_profile": focus_profile
        },
        "focus_profile": focus_profile,
        "conclusion": "evidence-backed company_profile bundle generated from unified company profile input.",
        "synthesis_mode": synth.get("mode"),
        "sections": sections,
        "artifact_plan": [
            {"artifact_type": "word", "path": str(docx_path.relative_to(ROOT))},
            {"artifact_type": "excel", "path": str(xlsx_path.relative_to(ROOT))},
            {"artifact_type": "ppt", "path": str(pptx_path.relative_to(ROOT))},
            {"artifact_type": "json", "path": str(json_path.relative_to(ROOT))}
        ]
    }

    write_json(json_path, bundle)
    print(json.dumps(bundle, ensure_ascii=False, indent=2))

if __name__ == "__main__":
    main()
