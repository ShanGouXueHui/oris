#!/usr/bin/env python3
import argparse
import json
from pathlib import Path

from docx import Document
from openpyxl import load_workbook
from openpyxl.workbook import Workbook
from pptx import Presentation

ROOT = Path(__file__).resolve().parents[2]

SECTION_TITLE_MAP = {
    "executive_summary": "Executive Summary",
    "company_or_case_overview": "Company / Case Overview",
    "industry_and_competition": "Industry and Competition",
    "technology_stack_breakdown": "Technology Stack Breakdown",
    "customer_scenario_analysis": "Customer Scenario Analysis",
    "recommendations": "Recommendations",
    "risks": "Risks",
    "citations_appendix": "Citations Appendix",
}

PPT_TITLE_MAP = {
    "title": "Title",
    "executive_summary": "Executive Summary",
    "industry_context": "Industry Context",
    "competitive_position": "Competitive Position",
    "technology_stack": "Technology Stack",
    "customer_scenarios": "Customer Scenarios",
    "recommendations": "Recommendations",
    "next_steps": "Next Steps",
}

def load_json(path: Path):
    return json.loads(path.read_text(encoding="utf-8"))

def unwrap_source_payload(obj: dict):
    if not isinstance(obj, dict):
        return {}
    if (obj.get("request") or {}).get("analysis_type") == "account_strategy":
        return obj
    if obj.get("analysis_type") == "account_strategy":
        return obj
    if obj.get("company_name") or obj.get("company_profile") or obj.get("core_data"):
        return obj
    return obj

def flatten_sources(case_obj: dict):
    rows = []
    seen = set()

    def push(source_name, source_type, url, publisher):
        key = (source_name or "", url or "")
        if key in seen:
            return
        seen.add(key)
        rows.append({
            "source_name": source_name or "",
            "source_type": source_type or "",
            "url": url or "",
            "publisher": publisher or ""
        })

    for ent in (case_obj.get("detected_entities") or []):
        for s in (ent.get("sources") or []):
            push(s.get("source_name"), s.get("source_type"), s.get("url"), s.get("publisher"))

    bindings = case_obj.get("role_bindings") or {}
    for key in ["partner", "cloud_vendor", "target_company"]:
        ent = bindings.get(key) or {}
        for s in (ent.get("sources") or []):
            push(s.get("source_name"), s.get("source_type"), s.get("url"), s.get("publisher"))
    for key in ["customers", "competitors"]:
        for ent in (bindings.get(key) or []):
            for s in (ent.get("sources") or []):
                push(s.get("source_name"), s.get("source_type"), s.get("url"), s.get("publisher"))
    return rows

def flatten_citations(source_obj: dict):
    rows = []
    seen = set()

    dbb = source_obj.get("db_backed_benchmark") or {}
    for entity_name, block in dbb.items():
        for c in (block.get("recent_citations") or []):
            key = (c.get("citation_label"), c.get("citation_url"))
            if key in seen:
                continue
            seen.add(key)
            rows.append({
                "entity_name": entity_name,
                "citation_label": c.get("citation_label") or "",
                "citation_url": c.get("citation_url") or "",
                "claim_code": c.get("claim_code") or "",
                "citation_note": c.get("citation_note") or ""
            })
    return rows

def flatten_claim_usage(source_obj: dict):
    rows = []
    for q in (source_obj.get("questions") or []):
        rows.append({"claim_group": "question", "claim_text": q})
    for t in (source_obj.get("recommendation_titles") or []):
        rows.append({"claim_group": "recommendation_title", "claim_text": t})
    for r in (source_obj.get("risks") or []):
        rows.append({"claim_group": "risk", "claim_text": r})
    return rows

def flatten_competitor_matrix(source_obj: dict):
    rows = []
    for row in (source_obj.get("competitor_matrix") or []):
        rows.append({
            "entity_name": row.get("entity_name") or "",
            "entity_type": row.get("entity_type") or "",
            "signal_strength": row.get("signal_strength"),
            "source_snapshot_count_recent": row.get("source_snapshot_count_recent"),
            "evidence_item_count_recent": row.get("evidence_item_count_recent"),
            "metric_observation_count_recent": row.get("metric_observation_count_recent"),
            "citation_count_recent": row.get("citation_count_recent"),
            "dimensions": ", ".join(row.get("dimensions") or [])
        })
    return rows

def evolution_actions_by_type(compiled_case: dict):
    out = {}
    for item in (compiled_case.get("evolution_actions") or []):
        out.setdefault(item.get("action_type"), []).append(item.get("value"))
    return out

def ensure_docx_sections(docx_path: Path, source_obj: dict, compiled_case: dict):
    doc = Document(str(docx_path))
    actions = evolution_actions_by_type(compiled_case)

    word_sections = actions.get("upgrade_word_structure_candidate", [])
    if not word_sections:
        doc.save(str(docx_path))
        return {"updated": False, "added_sections": 0}

    matrix = flatten_competitor_matrix(source_obj)
    sources = flatten_sources(compiled_case)
    citations = flatten_citations(source_obj)
    questions = source_obj.get("questions") or []
    recs = source_obj.get("recommendation_titles") or []
    risks = source_obj.get("risks") or []

    for sec in word_sections:
        title = SECTION_TITLE_MAP.get(sec, sec.replace("_", " ").title())
        doc.add_heading(title, level=1)

        if sec == "company_or_case_overview":
            req = (source_obj.get("request") or {})
            partner = ((req.get("partner") or {}).get("name")) or ""
            cloud = ((req.get("cloud_vendor") or {}).get("name")) or ""
            customers = ", ".join([x.get("name", "") for x in (req.get("customers") or []) if x.get("name")])
            competitors = ", ".join([x.get("name", "") for x in (req.get("competitors") or []) if x.get("name")])
            doc.add_paragraph(f"Partner: {partner}")
            doc.add_paragraph(f"Cloud Vendor: {cloud}")
            doc.add_paragraph(f"Customers: {customers}")
            doc.add_paragraph(f"Competitors: {competitors}")

        elif sec == "industry_and_competition":
            if matrix:
                for row in matrix[:8]:
                    doc.add_paragraph(
                        f'{row["entity_name"]}: signal_strength={row["signal_strength"]}, '
                        f'snapshots={row["source_snapshot_count_recent"]}, '
                        f'evidence={row["evidence_item_count_recent"]}, '
                        f'citations={row["citation_count_recent"]}'
                    )
            else:
                doc.add_paragraph("No competitor matrix was found in source JSON.")

        elif sec == "technology_stack_breakdown":
            cap_map = source_obj.get("capability_mapping") or []
            if cap_map:
                for row in cap_map:
                    doc.add_paragraph(
                        f'{row.get("mapping_code")}: '
                        f'partner={row.get("partner_strength")} | '
                        f'cloud={row.get("cloud_strength")} | '
                        f'value={row.get("joint_value")}'
                    )
            else:
                dims = ", ".join(source_obj.get("dimensions") or compiled_case.get("dimensions") or [])
                doc.add_paragraph(f"Technology-related focus dimensions: {dims}")

        elif sec == "customer_scenario_analysis":
            for c in ((source_obj.get("request") or {}).get("customers") or []):
                doc.add_paragraph(f'Customer scenario: {c.get("name")} ({c.get("type") or c.get("role") or ""})')
            if questions:
                for q in questions[:6]:
                    doc.add_paragraph(f'Question: {q}')

        elif sec == "recommendations":
            if recs:
                for x in recs:
                    doc.add_paragraph(x)
            else:
                doc.add_paragraph("No recommendation titles were found.")

        elif sec == "risks":
            if risks:
                for x in risks:
                    doc.add_paragraph(x)
            else:
                comp = compiled_case.get("compare_summary") or {}
                doc.add_paragraph(f'Compare summary: {json.dumps(comp, ensure_ascii=False)}')

        elif sec == "citations_appendix":
            for row in citations[:40]:
                doc.add_paragraph(f'{row["entity_name"]} | {row["citation_label"]} | {row["citation_url"]}')
            if not citations:
                for row in sources[:20]:
                    doc.add_paragraph(f'{row["source_name"]} | {row["url"]}')

        else:
            doc.add_paragraph(f"Generated by ORIS evolution pipeline for section: {sec}")

    doc.save(str(docx_path))
    return {"updated": True, "added_sections": len(word_sections)}

def ensure_xlsx_sheets(xlsx_path: Path, source_obj: dict, compiled_case: dict):
    try:
        wb = load_workbook(str(xlsx_path))
    except Exception:
        wb = Workbook()

    actions = evolution_actions_by_type(compiled_case)
    excel_sections = actions.get("upgrade_excel_evidence_candidate", [])

    def replace_sheet(title, headers, rows):
        if title in wb.sheetnames:
            ws = wb[title]
            wb.remove(ws)
        ws = wb.create_sheet(title)
        ws.append(headers)
        for row in rows:
            ws.append([row.get(h, "") for h in headers])

    if "raw_evidence_table" in excel_sections:
        rows = flatten_competitor_matrix(source_obj)
        replace_sheet("RawEvidence", list(rows[0].keys()) if rows else ["entity_name", "signal_strength"], rows)

    if "source_registry" in excel_sections:
        rows = flatten_sources(compiled_case)
        replace_sheet("SourceRegistry", list(rows[0].keys()) if rows else ["source_name", "source_type", "url", "publisher"], rows)

    if "citation_binding_table" in excel_sections:
        rows = flatten_citations(source_obj)
        replace_sheet("CitationBinding", list(rows[0].keys()) if rows else ["entity_name", "citation_label", "citation_url", "claim_code", "citation_note"], rows)

    if "claim_usage_table" in excel_sections:
        rows = flatten_claim_usage(source_obj)
        replace_sheet("ClaimUsage", list(rows[0].keys()) if rows else ["claim_group", "claim_text"], rows)

    if "Sheet" in wb.sheetnames and len(wb.sheetnames) > 1:
        ws = wb["Sheet"]
        if ws.max_row == 1 and ws.max_column == 1 and ws["A1"].value is None:
            wb.remove(ws)

    wb.save(str(xlsx_path))
    return {"updated": True, "added_sheets": excel_sections}

def add_ppt_bullet_slide(prs, title, bullets):
    layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = None
    for ph in slide.placeholders:
        if getattr(ph, "placeholder_format", None) and ph.placeholder_format.idx != 0:
            body = ph
            break
    if body is None:
        tx = slide.shapes.add_textbox(left=600000, top=1400000, width=8000000, height=4500000)
        tf = tx.text_frame
    else:
        tf = body.text_frame
        tf.clear()
    if not bullets:
        bullets = ["No content available."]
    first = True
    for b in bullets[:12]:
        if first:
            tf.text = str(b)
            first = False
        else:
            p = tf.add_paragraph()
            p.text = str(b)

def ensure_ppt_slides(pptx_path: Path, source_obj: dict, compiled_case: dict):
    prs = Presentation(str(pptx_path))
    actions = evolution_actions_by_type(compiled_case)
    ppt_sections = actions.get("upgrade_ppt_story_candidate", [])

    matrix = flatten_competitor_matrix(source_obj)
    recs = source_obj.get("recommendation_titles") or []
    questions = source_obj.get("questions") or []
    cap_map = source_obj.get("capability_mapping") or []

    for sec in ppt_sections:
        title = PPT_TITLE_MAP.get(sec, sec.replace("_", " ").title())
        bullets = []

        if sec == "industry_context":
            bullets = [
                f'Questions: {len(questions)}',
                f'Competitors in matrix: {len(matrix)}',
                f'Capabilities mapped: {len(cap_map)}'
            ]
        elif sec == "competitive_position":
            bullets = [
                f'{x["entity_name"]}: strength={x["signal_strength"]}, evidence={x["evidence_item_count_recent"]}, citations={x["citation_count_recent"]}'
                for x in matrix[:8]
            ]
        elif sec == "technology_stack":
            bullets = [
                f'{x.get("mapping_code")}: {x.get("joint_value")}'
                for x in cap_map[:8]
            ] or ["Technology stack detail not found in source JSON."]
        elif sec == "customer_scenarios":
            bullets = [f'Question: {q}' for q in questions[:8]]
        elif sec == "recommendations":
            bullets = recs[:8] or ["No recommendation titles found."]
        elif sec == "next_steps":
            bullets = [
                "Validate source coverage and citation density",
                "Expand technology-stack evidence",
                "Strengthen customer scenario and recommendation logic"
            ]
        else:
            bullets = [f"Generated by ORIS evolution pipeline for slide: {sec}"]

        add_ppt_bullet_slide(prs, title, bullets)

    prs.save(str(pptx_path))
    return {"updated": True, "added_slides": len(ppt_sections)}

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--analysis-type", required=True)
    ap.add_argument("--source-json-path", required=True)
    ap.add_argument("--compiled-case-path", required=True)
    ap.add_argument("--report-build-dir", required=True)
    args = ap.parse_args()

    source_obj = unwrap_source_payload(load_json(Path(args.source_json_path)))
    compiled_case = load_json(Path(args.compiled_case_path))
    report_dir = Path(args.report_build_dir)

    changed = {
        "docx": None,
        "xlsx": None,
        "pptx": None
    }

    for p in sorted(report_dir.glob("*.docx")):
        changed["docx"] = ensure_docx_sections(p, source_obj, compiled_case)

    for p in sorted(report_dir.glob("*.xlsx")):
        changed["xlsx"] = ensure_xlsx_sheets(p, source_obj, compiled_case)

    for p in sorted(report_dir.glob("*.pptx")):
        changed["pptx"] = ensure_ppt_slides(p, source_obj, compiled_case)

    print(json.dumps({
        "ok": True,
        "analysis_type": args.analysis_type,
        "report_build_dir": str(report_dir),
        "changed": changed
    }, ensure_ascii=False, indent=2))

if __name__ == "__main__":
    main()
